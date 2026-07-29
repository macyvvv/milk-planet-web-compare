import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "milk-planet Webサイト\nリニューアル提案"
subtitle.text = "現行サイトと新提案(branch5)の比較と改修点まとめ"

# Slide 2: 目的・コンセプト
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "目的とコンセプト"
body = slide.placeholders[1]
tf = body.text_frame
tf.text = "目的"
p = tf.add_paragraph()
p.text = "既存の画像素材・HTML構造を極力活かしつつ、使いやすさ(UI/UX)とモダンな見栄えを向上させ、運用コストを抑えた新サイトへ移行する。"
p.level = 1

p2 = tf.add_paragraph()
p2.text = "コンセプト"
p2.level = 0
p3 = tf.add_paragraph()
p3.text = "1. スマホ最適化 (モバイルファースト)"
p3.level = 1
p4 = tf.add_paragraph()
p4.text = "2. 導線の明確化 (ユーザーが迷わない設計)"
p4.level = 1
p5 = tf.add_paragraph()
p5.text = "3. デザインの統一感 (ブランドイメージの向上)"
p5.level = 1

# Slide 3: 比較表
slide_layout = prs.slide_layouts[5] # blank slide with title
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "新旧サイト 比較表"

rows, cols = 6, 3
left = Inches(0.5)
top = Inches(1.5)
width = Inches(9.0)
height = Inches(3.5)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(3.5)
table.columns[2].width = Inches(3.5)

# Headers
headers = ["項目", "現行サイト (Current)", "新提案サイト (branch5)"]
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    cell.text_frame.paragraphs[0].font.bold = True

data = [
    ["メインビジュアル", "画像が見切れる場合がある", "元画像を切らずに最大表示。周辺画像との関係性も整理"],
    ["ナビゲーション", "ページ上部に常設（スマホで領域を圧迫）", "ハンバーガーメニューに統一し、コンテキストに合わせて配置"],
    ["導線・リンク", "空アンカー等により遷移先がブレる箇所がある", "遷移先を安定化。ハッシュやアンカーの誤動作を防止"],
    ["店舗ページ", "ページごとのUI・サイズ感が不揃い", "バナー、アクセス、ロゴのサイズ感を見直し、統一感を確保"],
    ["キャストページ", "余白が広く、SNSリンクが目立たない", "縦横の間隔を最適化し、SNSリンクを復元・強調"]
]

for row_idx, row_data in enumerate(data):
    for col_idx, cell_data in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = cell_data
        cell.text_frame.paragraphs[0].font.size = Pt(14)

# Slide 4: 主要な改修点 1
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "主要な改修点 1: トップページと導線"
body = slide.placeholders[1]
tf = body.text_frame
tf.text = "カルーセルの改善"
p = tf.add_paragraph()
p.text = "画面幅を基準にし、中央のメイン画像を主役に据える構成に変更。"
p.level = 1
p = tf.add_paragraph()
p.text = "前後の画像も少し見せる(1:8:1の比率)ことで、現在位置と全体の繋がりをわかりやすく表現。"
p.level = 1

p = tf.add_paragraph()
p.text = "ハンバーガーメニューの導入"
p.level = 0
p = tf.add_paragraph()
p.text = "スマホでの操作性を考慮し、全ページ共通で操作しやすいメニューボタンを配置。"
p.level = 1
p = tf.add_paragraph()
p.text = "「しすてむ＆めにゅう」など、必要な情報へのアクセスをシンプル化。"
p.level = 1

# Slide 5: 主要な改修点 2
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "主要な改修点 2: 店舗＆キャストページ"
body = slide.placeholders[1]
tf = body.text_frame
tf.text = "UIの統一と整理"
p = tf.add_paragraph()
p.text = "店舗一覧ページと個別店舗ページの見え方を統一。"
p.level = 1
p = tf.add_paragraph()
p.text = "バナーやアクセスマップの配置を調整し、オリジナルデザインの良さを保ちつつ洗練された印象に。"
p.level = 1

p = tf.add_paragraph()
p.text = "キャスト表示の改善"
p.level = 0
p = tf.add_paragraph()
p.text = "写真配置のズレを解消し、SNSリンクへの導線を復元。"
p.level = 1
p = tf.add_paragraph()
p.text = "既存のフィルター機能などのリッチな挙動は維持しつつ、使い勝手を向上。"
p.level = 1

# Slide 6: 実装・運用のメリット
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "実装・運用のメリット"
body = slide.placeholders[1]
tf = body.text_frame
tf.text = "既存素材の徹底活用"
p = tf.add_paragraph()
p.text = "新たな画像素材の作成は不要。現在の更新フローとコストを維持したまま移行可能。"
p.level = 1

p = tf.add_paragraph()
p.text = "保守性の向上"
p.level = 0
p = tf.add_paragraph()
p.text = "CSS/JSを中心に調整を行っているため、将来的なコンテンツ追加時もレイアウト崩れが起きにくい設計。"
p.level = 1

p = tf.add_paragraph()
p.text = "今後の拡張性"
p.level = 0
p = tf.add_paragraph()
p.text = "このリニューアル(branch5)をベースに、PC版の最適化やさらなる機能追加をシームレスに行うことが可能。"
p.level = 1

prs.save('/Users/ariel/Documents/tools/#_amagi/milk-planet _web改修案/outputs/milk-planet-proposal-v2.pptx')
print("Saved pptx successfully")
